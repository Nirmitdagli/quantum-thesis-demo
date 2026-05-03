"""Build the Master's-thesis defence presentation for Nirmit Dagli.

Audience: CS students + professors with limited QC background.
Tone: professional but with friendly cartoon-style analogies (the
"restaurant kitchen" idea from the website index.html).

13 slides, ~30-40 min talk + Q&A.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# -------------------------------------------------------------------------
# Paths
# -------------------------------------------------------------------------
ROOT      = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PLOTS_DIR = os.path.join(ROOT, "hybrid_simulation", "output", "plots")
OUT       = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                          "Defense_Presentation_Nirmit_Dagli.pptx")

def img(name):
    p = os.path.join(PLOTS_DIR, name)
    return p if os.path.exists(p) else None

# -------------------------------------------------------------------------
# Brand colours (match the website)
# -------------------------------------------------------------------------
BG       = RGBColor(0x0A, 0x0A, 0x1A)
PANEL    = RGBColor(0x14, 0x1E, 0x33)
TEXT     = RGBColor(0xE0, 0xE0, 0xF0)
MUTED    = RGBColor(0xA0, 0xA0, 0xC0)
BLUE     = RGBColor(0x4C, 0x9B, 0xE8)   # CPU
GREEN    = RGBColor(0x50, 0xC8, 0x78)   # GPU
PURPLE   = RGBColor(0x9C, 0x6A, 0xDE)   # QPU
ORANGE   = RGBColor(0xF0, 0x96, 0x3C)   # warnings
GOLD     = RGBColor(0xFF, 0xD7, 0x00)   # highlight
CYAN     = RGBColor(0x4D, 0xD9, 0xD9)
RED      = RGBColor(0xE8, 0x55, 0x5A)

# -------------------------------------------------------------------------
# Presentation setup — 16:9 widescreen
# -------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

# -------------------------------------------------------------------------
# Helpers
# -------------------------------------------------------------------------
def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, italic=False,
             color=TEXT, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=16, color=TEXT,
                level_color_map=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, lvl = item
        else:
            text, lvl = item, 0
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = lvl
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        bullet = "  " * lvl + ("• " if lvl == 0 else "– ")
        r.text = bullet + text
        r.font.size = Pt(size - lvl * 2)
        r.font.color.rgb = color if not (level_color_map and lvl in level_color_map) else level_color_map[lvl]
        r.font.name = "Calibri"
    return tb

def add_header(slide, slide_num, total, title, subtitle=None):
    """Standard slide header: thin bar with slide number, then big title."""
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  0, 0, prs.slide_width, Inches(0.45))
    bar.fill.solid(); bar.fill.fore_color.rgb = PANEL
    bar.line.fill.background()
    add_text(slide, 0.4, 0.07, 4, 0.3,
             f"HERO  ·  Defense  ·  {slide_num}/{total}",
             size=10, color=MUTED)
    add_text(slide, 9.5, 0.07, 3.5, 0.3,
             "Nirmit Dagli  ·  Quinnipiac · 2026",
             size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    # title
    add_text(slide, 0.5, 0.65, 12.3, 0.7, title,
             size=32, bold=True, color=TEXT)
    if subtitle:
        add_text(slide, 0.5, 1.35, 12.3, 0.4, subtitle,
                 size=16, italic=True, color=MUTED)

def add_image(slide, x, y, w, h, path):
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y),
                                  width=Inches(w), height=Inches(h))

def add_box(slide, x, y, w, h, color=PANEL, alpha=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.color.rgb = MUTED
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    return box

TOTAL = 13

# =========================================================================
# SLIDE 1 — TITLE
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
# Decorative purple accent bar
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4),
                          prs.slide_width, Inches(0.06))
acc.fill.solid(); acc.fill.fore_color.rgb = PURPLE
acc.line.fill.background(); acc.shadow.inherit = False

add_text(s, 0.6, 0.5, 12.1, 0.5, "MASTER'S THESIS DEFENCE  ·  MAY 2026",
         size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

add_text(s, 0.6, 0.95, 12.1, 1.2,
         "HERO",
         size=66, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_text(s, 0.6, 2.55, 12.1, 1.0,
         "Hybrid Quantum-AI Models for Cybersecurity",
         size=34, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 3.15, 12.1, 0.6,
         "on Cloud Platforms",
         size=34, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 3.95, 12.1, 0.5,
         "Accuracy  ·  Latency  ·  Energy",
         size=20, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Author & supervisors block
add_box(s, 3.0, 4.95, 7.3, 1.85)
add_text(s, 3.2, 5.10, 7.0, 0.5, "Nirmit Dagli",
         size=24, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_text(s, 3.2, 5.55, 7.0, 0.4,
         "Master of Science in Computer Science",
         size=14, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, 3.2, 5.95, 7.0, 0.4,
         "Quinnipiac University  ·  Hamden, Connecticut",
         size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, 3.2, 6.30, 7.0, 0.4,
         "Advisor: Prof. Taskin   ·   Committee: Prof. Kruti, Prof. Lin",
         size=13, color=BLUE, align=PP_ALIGN.CENTER)

add_text(s, 0.6, 7.05, 12.1, 0.3,
         "github.com/Nirmitdagli/quantum-thesis-demo   ·   nirmitdagli.github.io/quantum-thesis-demo/",
         size=10, color=MUTED, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 2 — THE PROBLEM
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 2, TOTAL, "Why Are We Here?",
           "Cybersecurity is exploding. Quantum is hyped. Cloud bills are real.")

# Three problem cards
cards = [
    ("Cybersecurity scale", BLUE,
     "Networks generate billions of\nsecurity events per day.\nClassical AI is keeping up — for now."),
    ("Quantum hype", PURPLE,
     "\"Quantum will revolutionise security!\"\nReally? At what cost?\nIn how many years?"),
    ("Cloud costs are real", ORANGE,
     "AWS Braket: $0.30/task on QPU.\nAWS EC2: $0.71/hour.\nWhich tier should you pay for?"),
]
for i, (title, col, text) in enumerate(cards):
    x = 0.5 + i * 4.27
    add_box(s, x, 1.95, 4.0, 2.6, col)
    add_text(s, x + 0.2, 2.10, 3.6, 0.5, title,
             size=20, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 2.65, 3.6, 1.8, text,
             size=14, color=TEXT, align=PP_ALIGN.CENTER)

add_box(s, 0.5, 4.85, 12.3, 1.8, PANEL)
add_text(s, 0.7, 5.0, 11.9, 0.5,
         "The question nobody has answered:",
         size=18, bold=True, color=GOLD)
add_text(s, 0.7, 5.5, 11.9, 1.1,
         'Given a cybersecurity workload, on which cloud tier should it run\n'
         '(CPU, GPU, or QPU) to minimise cost + energy + latency together — without losing accuracy?',
         size=20, italic=True, color=TEXT, align=PP_ALIGN.LEFT)

add_text(s, 0.5, 7.0, 12.3, 0.4,
         "This thesis answers that question with measurement, not opinion.",
         size=14, italic=True, color=BLUE, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 3 — THE RESTAURANT KITCHEN ANALOGY
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 3, TOTAL, "Think of It Like a Restaurant Kitchen",
           "How CPU + GPU + QPU divide the work")

# Three chef cards
cards = [
    ("CPU\nThe Head Chef", BLUE,
     "Coordinates everything.\nReads orders, decides what to do.\nSmart but only one thing at a time.",
     "Data prep · Classical ML · Control"),
    ("GPU\n20 Sous-Chefs", GREEN,
     "Chops vegetables in parallel.\nAll 20 cooks at once.\nGreat for repetitive AI tasks.",
     "Neural networks · Parallel kernels"),
    ("QPU\nThe Magic Taste-Tester", PURPLE,
     "Tastes ALL ingredient combinations\nat once. Tells you the BEST one.\nMagical but expensive and rare.",
     "Quantum kernels · Optimisation · Search"),
]
for i, (title, col, body, footer) in enumerate(cards):
    x = 0.5 + i * 4.27
    add_box(s, x, 1.95, 4.0, 4.4, col)
    add_text(s, x + 0.2, 2.15, 3.6, 1.0, title,
             size=24, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 3.45, 3.6, 1.7, body,
             size=14, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 5.55, 3.6, 0.7, footer,
             size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

add_text(s, 0.5, 6.55, 12.3, 0.4,
         'You don\'t fire the head chef and replace them with a magic taste-tester.',
         size=15, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 6.95, 12.3, 0.4,
         "You use ALL THREE — each for what it does best.",
         size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 4 — INTRODUCING HERO
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 4, TOTAL, "HERO — A Measurement Framework",
           "Heterogeneous Energy-aware Runtime Orchestrator")

add_image(s, 0.5, 1.85, 7.0, 5.2, img("hybrid_architecture_live.png"))

# Right panel
add_box(s, 7.8, 1.95, 5.0, 5.0, PANEL)
add_text(s, 8.0, 2.10, 4.6, 0.5, "What HERO does",
         size=20, bold=True, color=GOLD)
add_bullets(s, 8.0, 2.65, 4.6, 4.2,
            ["Sends the same workload to all 3 tiers (CPU, GPU, QPU)",
             "Measures runtime, energy, accuracy per task",
             "Converts to 2026 cloud dollar cost",
             "Repeats 30× for statistical confidence",
             "Picks the Pareto-optimal tier",
             "Returns a recommendation"],
            size=14)

add_text(s, 0.5, 7.05, 12.3, 0.35,
         "Backend-agnostic — swap the simulator for real IBM/AWS/Azure QPU in one line.",
         size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 5 — ALGORITHM 1
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 5, TOTAL, "Algorithm 1 — One Page, Five Steps",
           "Allocation → Measurement → Aggregation → Pareto → Selection")

steps = [
    ("1", "ALLOCATE TIER", BLUE,
     "Look at the workload type and size.\n"
     "Use the rule book to pick a tier:\n"
     "molecular & big? → QPU.\n"
     "tabular & small? → CPU.\n"
     "Image? → GPU."),
    ("2", "MEASURE 30×", GREEN,
     "Run on every tier 30 times.\n"
     "Record: accuracy, latency,\n"
     "energy (joules), cost (USD)."),
    ("3", "AGGREGATE", CYAN,
     "Mean ± standard deviation.\n"
     "Wilcoxon p-values vs baseline.\n"
     "Box plots."),
    ("4", "PARETO", ORANGE,
     "Drop methods that are worse than\n"
     "another on every axis (cost,\n"
     "energy, latency, accuracy)."),
    ("5", "RECOMMEND", PURPLE,
     "Argmin over the Pareto set with\n"
     "the user's weighting of cost\n"
     "vs energy vs latency vs accuracy."),
]
for i, (n, title, col, text) in enumerate(steps):
    x = 0.4 + i * 2.55
    add_box(s, x, 1.95, 2.4, 4.7, col)
    # big number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x + 0.85), Inches(2.10),
                                 Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = TEXT
    circle.line.fill.background(); circle.shadow.inherit = False
    add_text(s, x + 0.85, 2.18, 0.7, 0.5, n,
             size=24, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 3.0, 2.2, 0.5, title,
             size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, 3.55, 2.1, 2.8, text,
             size=11, color=TEXT, align=PP_ALIGN.CENTER)

add_text(s, 0.5, 6.95, 12.3, 0.4,
         "Same algorithm, any workload. Same algorithm, any future hardware.",
         size=14, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 6 — TWO CONTRASTING WORKLOADS
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 6, TOTAL, "Two Workloads, Two Verdicts",
           "Quantum isn't always better. Quantum isn't always worse.")

# Left workload card: Cybersecurity
add_box(s, 0.5, 1.95, 6.0, 5.1, BLUE)
add_text(s, 0.7, 2.10, 5.6, 0.5, "Workload 1 — Cybersecurity",
         size=20, bold=True, color=TEXT)
add_text(s, 0.7, 2.6, 5.6, 0.4, "KDD Cup 99 intrusion detection",
         size=13, italic=True, color=GOLD)
add_bullets(s, 0.7, 3.1, 5.6, 3.7,
            ["494,021 real network connections",
             "22 attack types (DoS, scanning, botnets)",
             "Methods: SVM, RF, GBM, KNN, QSVM",
             "30-seed Wilcoxon validation",
             "Verdict: classical wins by 12,970×"],
            size=14)
add_text(s, 0.7, 6.55, 5.6, 0.4,
         "Reality check for quantum cybersecurity hype.",
         size=12, italic=True, color=TEXT)

# Right workload card: Molecular
add_box(s, 6.85, 1.95, 6.0, 5.1, PURPLE)
add_text(s, 7.05, 2.10, 5.6, 0.5, "Workload 2 — Molecular Simulation",
         size=20, bold=True, color=TEXT)
add_text(s, 7.05, 2.6, 5.6, 0.4, "VQE on the H₂ molecule",
         size=13, italic=True, color=GOLD)
add_bullets(s, 7.05, 3.1, 5.6, 3.7,
            ["2 qubits (parity-reduced)",
             "Hardware-efficient ansatz, COBYLA",
             "Classical exact diagonalisation as truth",
             "Achieves chemical accuracy (< 1.6 mHa)",
             "Verdict: classical wins for H₂…"],
            size=14)
add_text(s, 7.05, 6.55, 5.6, 0.4,
         "…but classical scales O(2ⁿ). Quantum wins above n=50.",
         size=12, italic=True, color=GOLD)

# =========================================================================
# SLIDE 7 — CYBERSECURITY RESULT (THE BIG NUMBER)
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 7, TOTAL, "Cybersecurity — Classical Wins by 12,970×",
           "30 runs · KDD Cup 99 · 5 classifiers head-to-head")

add_image(s, 0.4, 1.85, 8.5, 4.7, img("classifier_comparison.png"))

# Right panel — headline numbers
add_box(s, 9.1, 1.95, 3.85, 4.6, PANEL)
add_text(s, 9.3, 2.1, 3.5, 0.4, "Headline Numbers",
         size=16, bold=True, color=GOLD)

stats = [
    ("Random Forest", "F1 = 0.998", GREEN),
    ("QSVM",          "F1 = 0.667", PURPLE),
    ("",              "", TEXT),
    ("Energy (RF)",   "22 J/task",  GREEN),
    ("Energy (QSVM)", "8,315 J/task", PURPLE),
    ("",              "", TEXT),
    ("Cost RF/1M",    "$1.99 USD",  GREEN),
    ("Cost QSVM/1M",  "$2,100,000", PURPLE),
]
for i, (label, val, col) in enumerate(stats):
    y = 2.6 + i * 0.45
    if not label:
        # divider
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(9.4), Inches(y + 0.18),
                                   Inches(3.4), Emu(8000))
        line.fill.solid(); line.fill.fore_color.rgb = MUTED
        line.line.fill.background(); line.shadow.inherit = False
        continue
    add_text(s, 9.3, y, 1.8, 0.4, label, size=12, color=TEXT)
    add_text(s, 11.05, y, 1.8, 0.4, val,
             size=13, bold=True, color=col, align=PP_ALIGN.RIGHT)

add_text(s, 0.5, 6.85, 12.3, 0.4,
         "QSVM also recalls only 50% of attacks. Operationally unacceptable for IDS today.",
         size=14, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 8 — MOLECULAR RESULT
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 8, TOTAL, "Molecular Simulation — VQE Works, but Wait",
           "30 runs · H₂ ground state · VQE vs exact diagonalisation")

add_image(s, 0.4, 1.85, 8.5, 4.7, img("vqe_convergence.png"))

# Right panel
add_box(s, 9.1, 1.95, 3.85, 4.6, PANEL)
add_text(s, 9.3, 2.1, 3.5, 0.4, "What we measured",
         size=16, bold=True, color=GOLD)

add_text(s, 9.3, 2.6, 3.5, 0.4, "Exact (NumPy)",
         size=12, color=TEXT)
add_text(s, 9.3, 2.95, 3.5, 0.4, "−1.857275 Hartree",
         size=14, bold=True, color=GREEN)

add_text(s, 9.3, 3.45, 3.5, 0.4, "VQE estimate",
         size=12, color=TEXT)
add_text(s, 9.3, 3.80, 3.5, 0.4, "−1.857 ± 0.012 Ha",
         size=14, bold=True, color=PURPLE)

add_text(s, 9.3, 4.30, 3.5, 0.4, "Error vs exact",
         size=12, color=TEXT)
add_text(s, 9.3, 4.65, 3.5, 0.4, "0.007 Ha (chem. accurate)",
         size=14, bold=True, color=PURPLE)

add_text(s, 9.3, 5.15, 3.5, 0.4, "But for H₂ today…",
         size=12, italic=True, color=MUTED)
add_text(s, 9.3, 5.50, 3.5, 0.4, "Classical wins 700,000×",
         size=14, bold=True, color=ORANGE)

add_box(s, 0.5, 6.85, 12.3, 0.55, PANEL)
add_text(s, 0.7, 6.95, 11.9, 0.4,
         "Why VQE matters: classical needs 16 PB of RAM at 50 qubits. "
         "Quantum stays polynomial. The crossover IS coming.",
         size=14, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 9 — PARETO + ALLOCATION
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 9, TOTAL, "Pareto Frontier + Allocation Rule",
           "From measurements to a one-page decision rule")

add_image(s, 0.3, 1.85, 6.5, 5.0, img("pareto_frontier.png"))
add_image(s, 7.0, 1.85, 6.0, 5.0, img("allocation_decision_tree.png"))

add_text(s, 0.5, 6.95, 12.3, 0.4,
         "Pareto-optimal on cybersecurity: SVM (cheapest) and Random Forest (most accurate). "
         "Everything else is dominated.",
         size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 10 — CLOUD COST + CROSSOVER
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 10, TOTAL, "What Does It Cost in 2026 Dollars?",
           "AWS EC2  ·  AWS g5  ·  AWS Braket  ·  IBM Quantum  ·  Azure Quantum")

add_image(s, 0.3, 1.85, 6.5, 5.0, img("cloud_cost_comparison.png"))
add_image(s, 7.0, 1.85, 6.0, 5.0, img("crossover_projection.png"))

add_box(s, 0.5, 6.85, 12.3, 0.55, ORANGE)
add_text(s, 0.7, 6.95, 11.9, 0.4,
         "$1.99 (CPU) vs $2,100,000 (QPU) per million classifications. "
         "Classical wins TODAY. Crossover projected at ~10⁴ features with fault-tolerant hardware.",
         size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 11 — LIVE DEMO
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 11, TOTAL, "Live Demo — Try It Yourself",
           "Open-source artifact: 151 MB of code, data, figures, papers")

# 3 link cards
links = [
    ("Live Website", BLUE,
     "nirmitdagli.github.io/\nquantum-thesis-demo/",
     "Polished narrative\nsummary of the thesis"),
    ("Live Simulator", PURPLE,
     "nirmitdagli.github.io/\nquantum-thesis-demo/\nsimulator.html",
     "Configure a workload,\nsee HERO's decision"),
    ("Open Code Repo", GREEN,
     "github.com/Nirmitdagli/\nquantum-thesis-demo",
     "MIT licensed.\nFork it. Run it."),
]
for i, (title, col, url, desc) in enumerate(links):
    x = 0.5 + i * 4.27
    add_box(s, x, 1.95, 4.0, 3.7, col)
    add_text(s, x + 0.2, 2.10, 3.6, 0.5, title,
             size=22, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 2.7, 3.6, 1.4, url,
             size=13, color=TEXT, align=PP_ALIGN.CENTER, font="Consolas")
    add_text(s, x + 0.2, 4.55, 3.6, 0.9, desc,
             size=14, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

add_box(s, 0.5, 5.95, 12.3, 1.4, PANEL)
add_text(s, 0.7, 6.10, 11.9, 0.5, "Reproduce in 5 minutes",
         size=18, bold=True, color=GOLD)
add_text(s, 0.7, 6.55, 11.9, 0.4,
         "$  git clone https://github.com/Nirmitdagli/quantum-thesis-demo.git",
         size=13, color=TEXT, font="Consolas")
add_text(s, 0.7, 6.90, 11.9, 0.4,
         "$  python -m hybrid_simulation.run_hybrid",
         size=13, color=TEXT, font="Consolas")

# =========================================================================
# SLIDE 12 — LIMITATIONS & FUTURE WORK
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 12, TOTAL, "Limitations & Future Work",
           "What this thesis does NOT claim — and what's next")

# Two columns
add_text(s, 0.5, 1.85, 6.0, 0.45, "Limitations (honest)",
         size=20, bold=True, color=ORANGE)
add_bullets(s, 0.5, 2.4, 6.0, 4.2,
            ["No real quantum hardware (Qiskit Aer simulator only)",
             "Small quantum sample (120 points — kernel cost is O(n²))",
             "QPU energy is modelled, not measured (no facility telemetry)",
             "Only KDD Cup 99 for cross-tier comparison",
             "Only H₂ for VQE — LiH and BeH₂ supported but not run",
             "Cloud pricing is an April 2026 snapshot"],
            size=14)

add_text(s, 6.85, 1.85, 6.0, 0.45, "Future Work",
         size=20, bold=True, color=GREEN)
add_bullets(s, 6.85, 2.4, 6.0, 4.2,
            ["Migrate QPU tier to IBM Eagle / AWS Braket Rigetti",
             "Compare noiseless simulator vs real-hardware QSVM",
             "Add LiH, BeH₂, larger molecules to VQE",
             "RAPL / NVML live energy telemetry",
             "Kubernetes operator for production deployment",
             "Cloud billing API for live cost (vs static dict)"],
            size=14)

add_box(s, 0.5, 6.85, 12.3, 0.55, PANEL)
add_text(s, 0.7, 6.95, 11.9, 0.4,
         "The framework is designed to absorb every one of these as cloud QPU hardware matures.",
         size=14, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 13 — CONCLUSION + Q&A
# =========================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
add_header(s, 13, TOTAL, "Conclusion",
           "What you should remember")

# Three takeaways
takeaways = [
    ("It's a measurement framework",
     "HERO doesn't claim quantum advantage.\n"
     "It measures CPU + GPU + QPU honestly\n"
     "across cost, energy, latency, accuracy."),
    ("Two contrasting verdicts",
     "Cybersecurity: classical wins by 12,970×\n"
     "Molecular sim: classical wins TODAY,\n"
     "but loses scaling-wise above n = 50 qubits."),
    ("Workload-aware allocation rule",
     "Concrete decision tree from measurements.\n"
     "Open source, reproducible in 5 minutes,\n"
     "extensible to any new workload or backend."),
]
for i, (t, body) in enumerate(takeaways):
    x = 0.5 + i * 4.27
    add_box(s, x, 1.85, 4.0, 3.0, PANEL)
    # top accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(x), Inches(1.85),
                              Inches(4.0), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = [BLUE, PURPLE, GOLD][i]
    bar.line.fill.background(); bar.shadow.inherit = False
    add_text(s, x + 0.2, 2.05, 3.6, 0.6, t,
             size=18, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 2.75, 3.6, 1.9, body,
             size=14, color=TEXT, align=PP_ALIGN.CENTER)

# QA box
add_box(s, 0.5, 5.05, 12.3, 1.95, PURPLE)
add_text(s, 0.7, 5.20, 11.9, 0.6, "Thank you. Questions?",
         size=32, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_text(s, 0.7, 5.95, 11.9, 0.4,
         "github.com/Nirmitdagli/quantum-thesis-demo",
         size=16, color=TEXT, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 0.7, 6.40, 11.9, 0.4,
         "nirmitdagli.github.io/quantum-thesis-demo/",
         size=16, color=TEXT, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 0.5, 7.10, 12.3, 0.3,
         "Nirmit Dagli  ·  Quinnipiac University  ·  May 2026",
         size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# -------------------------------------------------------------------------
prs.save(OUT)
import os
size_kb = os.path.getsize(OUT) / 1024
print(f"SUCCESS: {OUT}")
print(f"Size:    {size_kb:.1f} KB")
print(f"Slides:  {len(prs.slides)}")
